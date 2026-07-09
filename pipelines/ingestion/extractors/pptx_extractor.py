import base64
import io
import logging

from PIL import Image
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from backend.providers.base import ChatMessage
from backend.providers.client import NeuroFlowClient
from backend.providers.router import RoutingCriteria

from .base import ExtractedPage

logger = logging.getLogger(__name__)




async def extract_pptx(file_path: str, client: NeuroFlowClient) -> list[ExtractedPage]:
    prs = Presentation(file_path)
    pages: list[ExtractedPage] = []

    criteria = RoutingCriteria(task_type="classification", require_vision=True)

    for i, slide in enumerate(prs.slides):
        slide_text = []
        notes = ""
        images = []

        # Extract text and images from shapes
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                slide_text.append(shape.text.strip())

            if getattr(shape, "shape_type", None) == MSO_SHAPE_TYPE.PICTURE:
                try:
                    images.append(shape.image.blob)
                except Exception as e:
                    logger.info(f"Failed to extract image from shape: {e}")

        # Extract speaker notes
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()

        # Process images with Vision LLM
        descriptions = []
        for img_blob in images:
            try:
                with Image.open(io.BytesIO(img_blob)) as img:
                    processed_img = img
                    max_size = 1024
                    if max(processed_img.size) > max_size:
                        ratio = max_size / max(processed_img.size)
                        new_size = (
                            int(processed_img.width * ratio),
                            int(processed_img.height * ratio),
                        )  # noqa: E501
                        processed_img = processed_img.resize(new_size, Image.Resampling.LANCZOS)  # type: ignore

                    buffered = io.BytesIO()
                    img_format = processed_img.format if processed_img.format else "JPEG"
                    if processed_img.mode != "RGB":
                        processed_img = processed_img.convert("RGB")  # type: ignore
                        img_format = "JPEG"

                    processed_img.save(buffered, format=img_format)
                    img_base64 = base64.b64encode(buffered.getvalue()).decode("utf-8")
                    mime_type = f"image/{img_format.lower()}"

                    messages = [
                        ChatMessage(
                            role="user",
                            content=[
                                {
                                    "type": "text",
                                    "text": "Provide a highly detailed description of this image or diagram from a presentation slide.",  # noqa: E501
                                },
                                {
                                    "type": "image_url",
                                    "image_url": {"url": f"data:{mime_type};base64,{img_base64}"},
                                },
                            ],
                        )
                    ]

                    generation = await client.chat(messages, criteria)
                    descriptions.append(generation.content.strip())
            except Exception as e:
                logger.info(f"Vision LLM failed for slide {i + 1}: {e}")
                descriptions.append(f"[Image description failed: {e}]")

        # Combine everything
        combined_content = []
        if slide_text:
            combined_content.append("\n".join(slide_text))

        if notes:
            combined_content.append(f"Speaker Notes:\n{notes}")

        if descriptions:
            combined_content.append("Images on this slide:\n" + "\n\n".join(descriptions))

        final_text = "\n\n".join(combined_content).strip()

        pages.append(
            ExtractedPage(
                page_number=i + 1,
                content=final_text if final_text else "[Empty Slide]",
                content_type="text",
                metadata={"has_notes": bool(notes), "image_count": len(images)},
            )
        )

    return pages
